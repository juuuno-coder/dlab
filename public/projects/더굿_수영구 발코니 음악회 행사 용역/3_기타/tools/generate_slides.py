#!/usr/bin/env python3
"""Generate Reveal.js HTML from a .pptx file.

Usage:
  python tools/generate_slides.py "path/to/your.pptx" output_dir

Produces:
  - output_dir/content.json
  - output_dir/assets/ (images extracted)
  - output_dir/reveal.html
"""
import sys
import os
import json
from pptx import Presentation


def safe_mkdir(path):
    os.makedirs(path, exist_ok=True)


def extract(pptx_path, out_dir):
    prs = Presentation(pptx_path)
    assets_dir = os.path.join(out_dir, 'assets')
    safe_mkdir(assets_dir)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        slide_images = []
        img_count = 0
        for shape in slide.shapes:
            # text
            if hasattr(shape, 'text') and shape.text and shape.text.strip():
                slide_texts.append(shape.text.strip())
            # images: many shapes exposing image via .image
            try:
                img = shape.image
            except Exception:
                img = None
            if img is not None:
                img_count += 1
                ext = img.ext
                img_name = f'slide{i}_img{img_count}.{ext}'
                img_path = os.path.join(assets_dir, img_name)
                with open(img_path, 'wb') as f:
                    f.write(img.blob)
                slide_images.append(os.path.join('assets', img_name))

        slides.append({'page': i, 'text': "\n\n".join(slide_texts), 'images': slide_images})

    # write JSON
    content_path = os.path.join(out_dir, 'content.json')
    with open(content_path, 'w', encoding='utf-8') as f:
        json.dump({'title': os.path.basename(pptx_path), 'slides': slides}, f, ensure_ascii=False, indent=2)

    # generate HTML
    generate_html(content_path, out_dir)


def generate_html(content_json_path, out_dir):
    with open(content_json_path, 'r', encoding='utf-8') as f:
        content = json.load(f)

    slides_html = []
    for s in content['slides']:
        parts = []
        if s.get('text'):
            parts.append(f"<div class=\"slide-text\">{escape_html(s['text']).replace('\n', '<br>')}</div>")
        for img in s.get('images', []):
            parts.append(f"<div class=\"slide-image\"><img src=\"{img}\"></div>")
        slide_section = f"<section>\n{''.join(parts)}\n</section>"
        slides_html.append(slide_section)

    html_template = HTML_TEMPLATE.replace('<!--TITLE-->', escape_html(content.get('title', 'Presentation')))
    html_template = html_template.replace('<!--SLIDES-->', '\n'.join(slides_html))

    out_html = os.path.join(out_dir, 'reveal.html')
    with open(out_html, 'w', encoding='utf-8') as f:
        f.write(html_template)

    print(f"Generated: {out_html}")


def escape_html(s):
    return (s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'))


HTML_TEMPLATE = """<!doctype html>
<html lang="ko">
  <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title><!--TITLE--></title>
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/reveal.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/theme/black.css">
    <style>
      .slide-text { font-family: 'Noto Sans KR', Arial, sans-serif; font-size: 24px; color:#111; max-width:960px; margin:0 auto; padding:20px }
      .slide-image img { max-width:960px; display:block; margin:10px auto }
    </style>
  </head>
  <body>
    <div class="reveal">
      <div class="slides">
<!--SLIDES-->
      </div>
    </div>
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/reveal.js"></script>
    <script>
      Reveal.initialize({
        hash: true,
        width: 1280,
        height: 720,
        slideNumber: true
      });
    </script>
  </body>
</html>
"""


def main():
    if len(sys.argv) < 3:
        print('Usage: python tools/generate_slides.py path/to/presentation.pptx out_dir')
        sys.exit(1)
    pptx_path = sys.argv[1]
    out_dir = sys.argv[2]
    safe_mkdir(out_dir)
    extract(pptx_path, out_dir)


if __name__ == '__main__':
    main()
