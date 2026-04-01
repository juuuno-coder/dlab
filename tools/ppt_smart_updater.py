"""
PPT 스마트 업데이터
RFP JSON을 기반으로 기존 PPT의 메타데이터를 지능적으로 업데이트
"""

import sys
import json
from pathlib import Path
from pptx import Presentation


class PPTSmartUpdater:
    def __init__(self, pptx_path, rfp_json_path):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(pptx_path))

        # RFP JSON 로드
        with open(rfp_json_path, 'r', encoding='utf-8') as f:
            self.rfp = json.load(f)

        print(f"[스마트 업데이트] PPT 로드: {self.pptx_path.name}")
        print(f"총 {len(self.prs.slides)} 슬라이드")
        print(f"RFP 로드: {Path(rfp_json_path).name}")

        # 업데이트 매핑 생성
        self.update_map = self._build_update_map()

    def _build_update_map(self):
        """
        RFP JSON에서 업데이트 매핑 생성
        기존 PPT의 메타데이터와 새 RFP 데이터를 매핑
        """
        update_map = {}

        # 1. 사업명 (가장 중요)
        if '사업명' in self.rfp:
            # 기존 제안서의 사업명을 찾아서 새 사업명으로 교체
            # 예: "2026년 제6회 부산 봄꽃 전시회" → 새 사업명
            update_map['사업명'] = {
                'target': self.rfp['사업명'],
                'context': ['제안서', '용역', '행사']  # 이 단어 근처에 있으면 사업명
            }

        # 2. 발주처
        if '발주처' in self.rfp and '기관명' in self.rfp['발주처']:
            update_map['발주기관'] = {
                'target': self.rfp['발주처']['기관명'],
                'context': ['발주', '주최', '주관']
            }

        # 3. 사업개요
        if '사업개요' in self.rfp:
            개요 = self.rfp['사업개요']

            # 예산
            if '예산' in 개요:
                update_map['예산'] = {
                    'target': 개요['예산'],
                    'context': ['원', '예산', '총사업비']
                }

            # 행사일정
            if '행사일정' in 개요:
                update_map['행사일정'] = {
                    'target': 개요['행사일정'],
                    'context': ['기간', '일정', '행사기간']
                }

            # 행사장소
            if '행사장소' in 개요:
                update_map['행사장소'] = {
                    'target': 개요['행사장소'],
                    'context': ['장소', '위치', '개최지']
                }

            # 주제
            if '주제' in 개요:
                update_map['주제'] = {
                    'target': 개요['주제'],
                    'context': ['주제', '슬로건', '테마']
                }

        return update_map

    def find_and_extract_current_values(self):
        """
        현재 PPT에서 업데이트할 값들을 추출
        (사용자가 확인할 수 있도록)
        """
        print("\n[현재 값 추출] PPT에서 기존 메타데이터 찾기")
        print("=" * 80)

        current_values = {}

        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                text = shape.text

                # 각 필드별로 현재 값 추출
                for field_name, field_data in self.update_map.items():
                    contexts = field_data.get('context', [])

                    # 컨텍스트 키워드가 포함되어 있으면 현재 값으로 추정
                    if any(ctx in text for ctx in contexts):
                        # 이 shape의 텍스트가 현재 값일 가능성 높음
                        if field_name not in current_values:
                            current_values[field_name] = []

                        current_values[field_name].append({
                            'slide': slide_num,
                            'text': text[:100],
                            'full_text': text
                        })

        # 출력
        for field_name, occurrences in current_values.items():
            print(f"\n[{field_name}] 발견: {len(occurrences)}건")
            for occ in occurrences[:3]:  # 처음 3개만 표시
                print(f"  슬라이드 {occ['slide']}: {occ['text']}...")

        return current_values

    def update_with_confirmation(self, auto_approve=False):
        """
        확인 후 업데이트
        """
        print("\n[업데이트 계획] 다음과 같이 업데이트됩니다:")
        print("=" * 80)

        # 현재 값 추출
        current_values = self.find_and_extract_current_values()

        print("\n[업데이트 매핑]")
        for field_name, field_data in self.update_map.items():
            target_value = field_data['target']
            print(f"  {field_name} → {target_value}")

        if not auto_approve:
            response = input("\n계속하시겠습니까? (y/n): ")
            if response.lower() != 'y':
                print("취소되었습니다.")
                return 0

        # 실제 업데이트 수행
        return self._perform_smart_update(current_values)

    def _perform_smart_update(self, current_values):
        """
        스마트 업데이트 수행
        """
        print("\n[업데이트 실행]")
        print("=" * 80)

        total_updated = 0

        for field_name, occurrences in current_values.items():
            if field_name not in self.update_map:
                continue

            target_value = self.update_map[field_name]['target']

            for occ in occurrences:
                slide_num = occ['slide']
                old_text = occ['full_text']

                # 해당 슬라이드에서 텍스트 교체
                slide = self.prs.slides[slide_num - 1]

                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame"):
                        continue

                    if shape.text == old_text:
                        # 이 shape의 텍스트를 업데이트
                        # (단순 치환이 아니라 스마트하게)
                        updated = self._smart_replace_in_shape(
                            shape, field_name, target_value
                        )

                        if updated:
                            print(f"  슬라이드 {slide_num}: {field_name} 업데이트")
                            total_updated += 1

        print(f"\n[OK] 총 {total_updated}개 필드 업데이트 완료")
        return total_updated

    def _smart_replace_in_shape(self, shape, field_name, target_value):
        """
        Shape 내에서 스마트하게 텍스트 교체
        """
        contexts = self.update_map[field_name].get('context', [])

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text

                # 컨텍스트 키워드가 포함되어 있으면 교체
                if any(ctx in text for ctx in contexts):
                    # 기존 값을 찾아서 교체 (간단한 휴리스틱)
                    # 예: "예산: 110,000,000원" → "예산: 80,000,000원"
                    if ':' in text:
                        parts = text.split(':', 1)
                        if len(parts) == 2:
                            run.text = f"{parts[0]}: {target_value}"
                            return True
                    elif '：' in text:  # 전각 콜론
                        parts = text.split('：', 1)
                        if len(parts) == 2:
                            run.text = f"{parts[0]}：{target_value}"
                            return True

                # 숫자나 날짜 패턴 교체
                import re

                # 금액 패턴 (예: 110,000,000원)
                if field_name == '예산' and re.search(r'\d+(?:,\d+)*원', text):
                    run.text = re.sub(r'\d+(?:,\d+)*원', target_value, text)
                    return True

                # 날짜 패턴 (예: 2026.04.17)
                if field_name == '행사일정' and re.search(r'\d{4}\.\d{2}\.\d{2}', text):
                    run.text = re.sub(
                        r'\d{4}\.\d{2}\.\d{2}\s*~\s*\d{4}\.\d{2}\.\d{2}',
                        target_value,
                        text
                    )
                    return True

        return False

    def save(self, output_path=None):
        """저장"""
        if output_path is None:
            # 원본 파일 이름에 _updated 추가
            stem = self.pptx_path.stem
            suffix = self.pptx_path.suffix
            output_path = self.pptx_path.parent / f"{stem}_updated{suffix}"

        self.prs.save(str(output_path))
        print(f"\n[저장 완료] {output_path}")
        return output_path


def main():
    if len(sys.argv) < 3:
        print("사용법: python ppt_smart_updater.py <PPTX파일> <RFP_JSON>")
        print("\n예시:")
        print("  python ppt_smart_updater.py proposal.pptx rfp_analysis.json")
        print("\n이 도구는 RFP JSON의 메타데이터를 기반으로")
        print("기존 PPT의 사업명, 예산, 일정, 장소 등을 자동으로 업데이트합니다.")
        sys.exit(1)

    pptx_path = sys.argv[1]
    rfp_json_path = sys.argv[2]

    try:
        updater = PPTSmartUpdater(pptx_path, rfp_json_path)

        # 자동 승인 옵션
        auto_approve = '--yes' in sys.argv or '-y' in sys.argv

        # 업데이트 수행
        updated_count = updater.update_with_confirmation(auto_approve=auto_approve)

        if updated_count > 0:
            # 저장
            output_path = updater.save()

            print("\n" + "=" * 80)
            print("[완료] PPT 스마트 업데이트 완료!")
            print("=" * 80)
            print(f"\n업데이트된 파일: {output_path}")

    except Exception as e:
        print(f"\n[오류] 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
