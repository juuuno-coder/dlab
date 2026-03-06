"""
PPT 재현기
ppt_analyzer.py가 생성한 분석 JSON을 기반으로 동일한 PPT 재현
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTReconstructor:
    def __init__(self, analysis_json_path):
        with open(analysis_json_path, 'r', encoding='utf-8') as f:
            self.analysis = json.load(f)

        # 새 Presentation 생성
        self.prs = Presentation()

        # 슬라이드 크기 설정
        meta = self.analysis['메타정보']
        self.prs.slide_width = meta['슬라이드크기']['width']
        self.prs.slide_height = meta['슬라이드크기']['height']

        print(f"[재현] PPT 재현 시작: {meta['원본파일']}")
        print(f"총 {meta['총슬라이드수']} 슬라이드 재현 예정")

    def reconstruct_all(self):
        """전체 PPT 재현"""
        print("=" * 80)

        for slide_data in self.analysis['슬라이드']:
            slide_num = slide_data['슬라이드번호']
            print(f"슬라이드 {slide_num} 재현 중...", end='\r')
            self._reconstruct_slide(slide_data)

        print(f"\n[OK] 재현 완료! 총 {len(self.prs.slides)} 슬라이드")

    def _reconstruct_slide(self, slide_data):
        """단일 슬라이드 재현"""
        # 빈 레이아웃 사용 (모든 요소를 수동으로 추가)
        blank_layout = self.prs.slide_layouts[6]  # 일반적으로 6번이 빈 레이아웃
        slide = self.prs.slides.add_slide(blank_layout)

        # 요소별로 재현
        for element in slide_data['요소']:
            self._reconstruct_element(slide, element)

    def _reconstruct_element(self, slide, element):
        """단일 요소 재현"""
        element_type = element['타입']
        position = element['위치']

        # 텍스트박스 또는 플레이스홀더
        if element_type in ['텍스트박스', '플레이스홀더']:
            if '텍스트' in element:
                self._add_textbox(slide, element)

        # 표
        elif element_type == '표' and '표' in element:
            self._add_table(slide, element)

        # 이미지는 일단 스킵 (나중에 추가 가능)
        # elif element_type == '이미지':
        #     pass

    def _add_textbox(self, slide, element):
        """텍스트박스 추가"""
        position = element['위치']
        text_data = element['텍스트']

        # 텍스트박스 생성
        left = Emu(position['left'])
        top = Emu(position['top'])
        width = Emu(position['width'])
        height = Emu(position['height'])

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        # 단락별로 추가
        text_frame.clear()
        for para_idx, para_data in enumerate(text_data['단락']):
            if para_idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # 레벨 설정
            p.level = para_data.get('레벨', 0)

            # Run별로 텍스트 추가
            for run_idx, run_data in enumerate(para_data.get('runs', [])):
                if run_idx == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()

                run.text = run_data['텍스트']

                # 폰트 스타일 적용
                font_data = run_data.get('폰트', {})
                if font_data.get('이름'):
                    run.font.name = font_data['이름']
                if font_data.get('크기'):
                    run.font.size = Pt(font_data['크기'])
                if font_data.get('굵게'):
                    run.font.bold = True
                if font_data.get('기울임'):
                    run.font.italic = True

    def _add_table(self, slide, element):
        """표 추가"""
        position = element['위치']
        table_data = element['표']

        left = Emu(position['left'])
        top = Emu(position['top'])
        width = Emu(position['width'])
        height = Emu(position['height'])

        rows = table_data['행수']
        cols = table_data['열수']

        # 표 생성
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # 셀 데이터 채우기
        for cell_data in table_data['셀']:
            row_idx = cell_data['행']
            col_idx = cell_data['열']

            try:
                cell = table.rows[row_idx].cells[col_idx]
                cell.text = cell_data['텍스트']
            except IndexError:
                # 병합된 셀 등으로 인한 인덱스 오류 무시
                pass

    def save(self, output_path):
        """재현된 PPT 저장"""
        output_file = Path(output_path)
        self.prs.save(str(output_file))

        print(f"\n[OK] 재현 PPT 저장: {output_file}")
        print(f"   총 {len(self.prs.slides)} 슬라이드")

        return output_file


def main():
    if len(sys.argv) < 2:
        print("사용법: python ppt_reconstructor.py <분석JSON> [출력PPTX]")
        print("예시: python ppt_reconstructor.py analysis.json reconstructed.pptx")
        sys.exit(1)

    analysis_json = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else 'reconstructed.pptx'

    try:
        reconstructor = PPTReconstructor(analysis_json)
        reconstructor.reconstruct_all()
        reconstructor.save(output_pptx)

        print("\n" + "=" * 80)
        print("[완료] PPT 재현 완료!")
        print("=" * 80)
        print(f"\n다음 단계:")
        print(f"1. PowerPoint에서 {output_pptx} 열기")
        print(f"2. 원본과 비교하여 재현도 확인")
        print(f"3. 누락된 요소나 스타일 차이 점검")

    except Exception as e:
        print(f"\n[오류] 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
