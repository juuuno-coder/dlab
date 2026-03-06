"""
PPT 자동 생성기
Markdown 제안서 초안 + PPT 템플릿 JSON을 기반으로 PPTX 파일 자동 생성
"""

import sys
import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTGenerator:
    def __init__(self, markdown_path, template_json_path, base_pptx_path=None):
        # Markdown 제안서 로드
        self.md_path = Path(markdown_path)
        self.md_content = self.md_path.read_text(encoding='utf-8')

        # 템플릿 JSON 로드
        with open(template_json_path, 'r', encoding='utf-8') as f:
            self.template = json.load(f)

        # 베이스 PPTX가 있으면 로드, 없으면 새로 생성
        if base_pptx_path and Path(base_pptx_path).exists():
            self.prs = Presentation(str(base_pptx_path))
            print(f"[PPT] 베이스 PPT 로드: {base_pptx_path}")
        else:
            self.prs = Presentation()
            print(f"[PPT] 새 PPT 생성")

            # 슬라이드 크기 설정 (템플릿에서 가져오기)
            if 'metadata' in self.template:
                width = self.template['metadata'].get('slide_width', 10693400)
                height = self.template['metadata'].get('slide_height', 7562850)
                self.prs.slide_width = width
                self.prs.slide_height = height

    def parse_markdown(self):
        """Markdown 파싱하여 슬라이드별 내용 추출"""
        # 섹션 분리 (# 제목 기준)
        sections = re.split(r'\n# ', self.md_content)

        slides_data = []

        for section in sections:
            if not section.strip():
                continue

            # 제목과 본문 분리
            lines = section.split('\n')
            title = lines[0].strip()

            # 본문 (## 서브제목 기준으로 다시 분리)
            content = '\n'.join(lines[1:])

            # ## 서브제목이 있으면 슬라이드 분할
            subsections = re.split(r'\n## ', content)

            if len(subsections) > 1:
                # 첫 번째는 섹션 커버
                slides_data.append({
                    "type": "section_cover",
                    "title": title,
                    "content": subsections[0].strip()
                })

                # 나머지는 내용 슬라이드
                for subsection in subsections[1:]:
                    sub_lines = subsection.split('\n')
                    sub_title = sub_lines[0].strip()
                    sub_content = '\n'.join(sub_lines[1:]).strip()

                    slides_data.append({
                        "type": "content",
                        "title": sub_title,
                        "content": sub_content
                    })
            else:
                # 서브제목 없으면 단일 슬라이드
                slides_data.append({
                    "type": "content",
                    "title": title,
                    "content": content.strip()
                })

        return slides_data

    def generate_all_slides(self):
        """전체 슬라이드 생성"""
        slides_data = self.parse_markdown()

        print(f"[생성] 생성할 슬라이드: {len(slides_data)}개")
        print("=" * 80)

        for i, slide_data in enumerate(slides_data, 1):
            print(f"슬라이드 {i}/{len(slides_data)}: {slide_data['title'][:40]}...", end='\r')

            if slide_data['type'] == 'section_cover':
                self._add_section_cover_slide(slide_data)
            else:
                self._add_content_slide(slide_data)

        print(f"\n[OK] 슬라이드 생성 완료!")

    def _add_section_cover_slide(self, slide_data):
        """섹션 커버 슬라이드 추가"""
        # 레이아웃 1: 제목만 (섹션 구분용)
        slide_layout = self.prs.slide_layouts[1] if len(self.prs.slide_layouts) > 1 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # 제목
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            self._apply_title_style(slide.shapes.title)

    def _add_content_slide(self, slide_data):
        """내용 슬라이드 추가"""
        # 레이아웃 1: 제목 + 본문
        slide_layout = self.prs.slide_layouts[1] if len(self.prs.slide_layouts) > 1 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # 제목
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            self._apply_title_style(slide.shapes.title)

        # 본문
        content_text = slide_data['content']

        # 표가 있으면 표 추가
        if '|' in content_text and '---' in content_text:
            self._add_table_to_slide(slide, content_text)
        else:
            # 일반 텍스트
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                # 줄바꿈 기준으로 단락 추가
                lines = content_text.split('\n')
                for line in lines:
                    if line.strip():
                        # 번호 리스트 (1., -, * 등)
                        if re.match(r'^[\d\-\*]\.\s+', line.strip()):
                            p = tf.add_paragraph()
                            p.text = line.strip()
                            p.level = 0
                            self._apply_body_style(p)
                        # 일반 텍스트
                        else:
                            p = tf.add_paragraph()
                            p.text = line.strip()
                            p.level = 0
                            self._apply_body_style(p)

    def _add_table_to_slide(self, slide, table_text):
        """Markdown 표를 PPT 표로 변환"""
        # Markdown 표 파싱
        lines = [line.strip() for line in table_text.split('\n') if line.strip()]

        # 헤더와 데이터 분리
        header_line = None
        data_lines = []

        for line in lines:
            if line.startswith('|') and line.endswith('|'):
                if '---' in line:
                    continue
                if header_line is None:
                    header_line = line
                else:
                    data_lines.append(line)

        if not header_line:
            return

        # 셀 분리
        def parse_row(line):
            return [cell.strip() for cell in line.split('|')[1:-1]]

        header_cells = parse_row(header_line)
        rows_data = [parse_row(line) for line in data_lines]

        # PPT 표 추가
        rows = len(rows_data) + 1  # 헤더 포함
        cols = len(header_cells)

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 헤더
        for col_idx, cell_text in enumerate(header_cells):
            cell = table.rows[0].cells[col_idx]
            cell.text = cell_text
            # 헤더 스타일
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)

        # 데이터
        for row_idx, row_data in enumerate(rows_data, 1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.rows[row_idx].cells[col_idx]
                cell.text = cell_text
                # 본문 스타일
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)

    def _apply_title_style(self, shape):
        """제목 스타일 적용"""
        # 템플릿에서 폰트 정보 가져오기
        fonts = self.template.get('design_system', {}).get('fonts', [])

        # 기본 폰트 (템플릿에서 첫 번째 폰트)
        default_font = fonts[0] if fonts else 'Noto Sans KR'

        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = default_font
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.LEFT

    def _apply_body_style(self, paragraph):
        """본문 스타일 적용"""
        fonts = self.template.get('design_system', {}).get('fonts', [])
        default_font = fonts[1] if len(fonts) > 1 else fonts[0] if fonts else 'Noto Sans KR'

        paragraph.font.name = default_font
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.LEFT

    def save(self, output_path):
        """PPTX 파일 저장"""
        output_file = Path(output_path)
        self.prs.save(str(output_file))

        print(f"\n[OK] PPT 저장 완료: {output_file}")
        print(f"   총 {len(self.prs.slides)} 슬라이드")

        return output_file


def main():
    if len(sys.argv) < 3:
        print("사용법: python ppt_generator.py <Markdown> <템플릿JSON> [베이스PPTX] [출력PPTX]")
        print("예시: python ppt_generator.py proposal_draft.md template.json base.pptx output.pptx")
        print("\n베이스PPTX가 없으면 새로 생성됩니다.")
        sys.exit(1)

    md_path = sys.argv[1]
    template_json = sys.argv[2]
    base_pptx = sys.argv[3] if len(sys.argv) > 3 else None
    output_pptx = sys.argv[4] if len(sys.argv) > 4 else 'generated_proposal.pptx'

    try:
        generator = PPTGenerator(md_path, template_json, base_pptx)
        generator.generate_all_slides()
        generator.save(output_pptx)

        print("\n" + "=" * 80)
        print("[완료] 제안서 PPT 생성 완료!")
        print("=" * 80)
        print(f"\n다음 단계:")
        print(f"1. PowerPoint에서 {output_pptx} 열기")
        print(f"2. 디자인 및 레이아웃 조정")
        print(f"3. 이미지 및 차트 추가")
        print(f"4. 최종 검토 후 제출")

    except Exception as e:
        print(f"\n[오류] 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
