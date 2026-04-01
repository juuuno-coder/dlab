"""
PPT 편집 도구 - AI가 PPTX 파일을 직접 읽고 수정
python-pptx 기반
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN


class PPTEditor:
    def __init__(self, pptx_path):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(pptx_path))
        print(f"PPT 로드: {self.pptx_path.name}")
        print(f"슬라이드 수: {len(self.prs.slides)}")

    def find_text(self, search_text):
        """텍스트 찾기 (슬라이드 번호 + 위치 반환)"""
        results = []

        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                if search_text in shape.text:
                    results.append({
                        'slide_num': slide_num,
                        'shape_text': shape.text,
                        'contains': search_text
                    })

        return results

    def replace_text(self, old_text, new_text):
        """텍스트 치환 (모든 슬라이드)"""
        replaced_count = 0

        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                if old_text in shape.text:
                    # 텍스트 프레임이 있는 경우
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    replaced_count += 1
                                    print(f"  슬라이드 {slide_num}: '{old_text}' → '{new_text}'")

        return replaced_count

    def list_all_text(self):
        """모든 텍스트 출력 (슬라이드별)"""
        for slide_num, slide in enumerate(self.prs.slides, 1):
            print(f"\n슬라이드 {slide_num}:")
            print("-" * 80)

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"  • {shape.text[:100]}")

    def get_slide_info(self, slide_num):
        """특정 슬라이드 상세 정보"""
        if slide_num < 1 or slide_num > len(self.prs.slides):
            print(f"슬라이드 번호 오류: {slide_num}")
            return

        slide = self.prs.slides[slide_num - 1]

        print(f"\n슬라이드 {slide_num} 상세 정보:")
        print("=" * 80)

        for i, shape in enumerate(slide.shapes, 1):
            print(f"\n도형 {i}:")
            print(f"  타입: {shape.shape_type}")

            if hasattr(shape, "text"):
                print(f"  텍스트: {shape.text[:200]}")

            if hasattr(shape, "text_frame"):
                for para_num, para in enumerate(shape.text_frame.paragraphs, 1):
                    print(f"    단락 {para_num}:")
                    for run in para.runs:
                        print(f"      - {run.text[:100]}")
                        print(f"        폰트: {run.font.name}, 크기: {run.font.size}")

    def add_slide(self, title, content_lines):
        """새로운 슬라이드 추가 (제목 + 본문)"""
        # 레이아웃: 1 = 제목 + 본문
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title

        # 본문
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

        print(f"슬라이드 추가: '{title}'")

    def delete_slide(self, slide_num):
        """슬라이드 삭제"""
        if slide_num < 1 or slide_num > len(self.prs.slides):
            print(f"슬라이드 번호 오류: {slide_num}")
            return

        # python-pptx는 직접 삭제를 지원하지 않으므로 XML 조작 필요
        rId = self.prs.slides._sldIdLst[slide_num - 1].rId
        self.prs.part.drop_rel(rId)
        del self.prs.slides._sldIdLst[slide_num - 1]

        print(f"슬라이드 {slide_num} 삭제됨")

    def save(self, output_path=None):
        """저장"""
        if output_path is None:
            # 원본 파일 이름에 _edited 추가
            stem = self.pptx_path.stem
            suffix = self.pptx_path.suffix
            output_path = self.pptx_path.parent / f"{stem}_edited{suffix}"

        self.prs.save(str(output_path))
        print(f"\n저장 완료: {output_path}")
        return output_path


def main():
    if len(sys.argv) < 2:
        print("사용법:")
        print("  python ppt_editor.py <PPTX파일> --list              # 모든 텍스트 출력")
        print("  python ppt_editor.py <PPTX파일> --find '검색어'      # 텍스트 찾기")
        print("  python ppt_editor.py <PPTX파일> --replace '이전' '새로운'  # 치환")
        print("  python ppt_editor.py <PPTX파일> --info 3           # 슬라이드 3 상세 정보")
        sys.exit(1)

    pptx_path = sys.argv[1]
    editor = PPTEditor(pptx_path)

    if len(sys.argv) < 3:
        # 기본: 모든 텍스트 출력
        editor.list_all_text()
        return

    command = sys.argv[2]

    if command == "--list":
        editor.list_all_text()

    elif command == "--find" and len(sys.argv) > 3:
        search_text = sys.argv[3]
        results = editor.find_text(search_text)

        print(f"\n'{search_text}' 검색 결과: {len(results)}건")
        for result in results:
            print(f"  슬라이드 {result['slide_num']}: {result['shape_text'][:100]}")

    elif command == "--replace" and len(sys.argv) > 4:
        old_text = sys.argv[3]
        new_text = sys.argv[4]

        print(f"\n'{old_text}' → '{new_text}' 치환 중...")
        count = editor.replace_text(old_text, new_text)

        print(f"\n총 {count}개 치환됨")

        if count > 0:
            output_path = editor.save()
            print(f"수정된 파일: {output_path}")

    elif command == "--info" and len(sys.argv) > 3:
        slide_num = int(sys.argv[3])
        editor.get_slide_info(slide_num)

    else:
        print("잘못된 명령어")


if __name__ == '__main__':
    main()
