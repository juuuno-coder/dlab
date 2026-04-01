"""
제안서 내용 학습 도구
기존 PPT 제안서에서 섹션별 내용 패턴과 문장 스타일 학습
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from collections import defaultdict
import re


class ProposalContentLearner:
    def __init__(self, pptx_path, template_json_path=None):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(pptx_path))

        # 템플릿 JSON이 있으면 로드
        self.template_data = None
        if template_json_path:
            with open(template_json_path, 'r', encoding='utf-8') as f:
                self.template_data = json.load(f)

        self.content_patterns = {
            "섹션구조": [],
            "텍스트패턴": {},
            "강조표현": [],
            "숫자데이터": [],
            "성공사례": []
        }

    def learn_all(self):
        """전체 PPT에서 내용 패턴 학습"""
        print(f"[학습] 제안서 내용 학습 시작: {self.pptx_path.name}")
        print(f"총 {len(self.prs.slides)} 슬라이드")
        print("=" * 80)

        # 1. 섹션 구조 학습
        self._learn_section_structure()

        # 2. 슬라이드별 내용 패턴 학습
        for slide_num, slide in enumerate(self.prs.slides, 1):
            print(f"슬라이드 {slide_num} 학습 중...", end='\r')
            self._learn_slide_content(slide, slide_num)

        print(f"\n[OK] 학습 완료!")

        # 3. 학습 결과 요약
        self._summarize_learning()

        return self.content_patterns

    def _learn_section_structure(self):
        """섹션 구조 학습 (대제목 기준)"""
        sections = []
        current_section = None

        for slide_num, slide in enumerate(self.prs.slides, 1):
            # 제목이 있는 슬라이드만
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()

                # Part, 장, 절 등의 키워드로 대제목 판단
                if any(kw in title for kw in ['Part', 'PART', '장.', 'Ⅰ', 'Ⅱ', 'Ⅲ', 'Ⅳ', 'Ⅴ']):
                    # 중복 체크: 같은 제목의 섹션이 이미 있는지 확인
                    existing_section = next((s for s in sections if s['제목'] == title), None)

                    if existing_section:
                        # 기존 섹션을 current_section으로 설정 (하위 슬라이드 추가용)
                        current_section = existing_section
                    else:
                        # 새 섹션 생성
                        if current_section and current_section not in sections:
                            sections.append(current_section)

                        current_section = {
                            "제목": title,
                            "시작슬라이드": slide_num,
                            "하위슬라이드": []
                        }
                        sections.append(current_section)
                elif current_section:
                    current_section["하위슬라이드"].append({
                        "슬라이드번호": slide_num,
                        "제목": title
                    })

        self.content_patterns["섹션구조"] = sections

    def _learn_slide_content(self, slide, slide_num):
        """단일 슬라이드 내용 학습"""
        # 슬라이드 제목
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # 모든 텍스트 수집
        all_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_texts.append(shape.text.strip())

        # 패턴별 학습
        for text in all_texts:
            # 1. 강조 표현 학습
            self._learn_emphasis(text)

            # 2. 숫자 데이터 학습
            self._learn_numbers(text)

            # 3. 성공 사례 학습
            self._learn_success_cases(text)

        # 4. 제목별 텍스트 패턴 저장
        if title:
            if title not in self.content_patterns["텍스트패턴"]:
                self.content_patterns["텍스트패턴"][title] = []

            self.content_patterns["텍스트패턴"][title].extend(all_texts)

    def _learn_emphasis(self, text):
        """강조 표현 학습"""
        emphasis_keywords = [
            '차별화', '경쟁력', '핵심', '강점', '성공',
            '최고', '최적', '최선', '최상',
            '보장', '확보', '달성', '실현',
            '혁신', '창의', '독창',
            '탁월', '우수', '뛰어난'
        ]

        for keyword in emphasis_keywords:
            if keyword in text:
                # 문장 단위로 추출
                sentences = re.split(r'[.!?]\s+', text)
                for sentence in sentences:
                    if keyword in sentence and len(sentence) < 200:
                        if sentence not in self.content_patterns["강조표현"]:
                            self.content_patterns["강조표현"].append(sentence.strip())

    def _learn_numbers(self, text):
        """숫자 데이터 학습 (성과 지표 등)"""
        # 패턴: 숫자 + 단위
        number_patterns = [
            r'(\d+(?:,\d+)*)\s*(?:명|건|개|회|년|일|시간|분)',
            r'(\d+(?:\.\d+)?)\s*%',
            r'(\d+(?:,\d+)*)\s*(?:억|천만|만)?\s*원',
        ]

        for pattern in number_patterns:
            matches = re.findall(pattern, text)
            if matches:
                # 문맥과 함께 저장
                for match in matches:
                    context = self._get_context(text, str(match), 50)
                    if context not in self.content_patterns["숫자데이터"]:
                        self.content_patterns["숫자데이터"].append(context)

    def _learn_success_cases(self, text):
        """성공 사례 학습"""
        success_keywords = [
            '수행', '완료', '달성', '성공',
            '실적', '경험', '참여', '진행'
        ]

        for keyword in success_keywords:
            if keyword in text:
                # 문장 단위로 추출
                sentences = re.split(r'[.!?]\s+', text)
                for sentence in sentences:
                    if keyword in sentence and len(sentence) > 20 and len(sentence) < 300:
                        # 연도가 포함된 경우 (실제 사례일 가능성 높음)
                        if re.search(r'20\d{2}', sentence):
                            if sentence not in self.content_patterns["성공사례"]:
                                self.content_patterns["성공사례"].append(sentence.strip())

    def _get_context(self, text, keyword, context_length=50):
        """키워드 주변 문맥 추출"""
        idx = text.find(keyword)
        if idx == -1:
            return ""

        start = max(0, idx - context_length)
        end = min(len(text), idx + len(keyword) + context_length)

        context = text[start:end].strip()
        return context

    def _summarize_learning(self):
        """학습 결과 요약"""
        print("\n" + "=" * 80)
        print("[요약] 학습 결과 요약")
        print("=" * 80)

        print(f"\n[구조] 섹션 구조: {len(self.content_patterns['섹션구조'])}개 섹션")
        for section in self.content_patterns['섹션구조']:
            print(f"  - {section['제목']} (슬라이드 {section['시작슬라이드']}부터)")

        print(f"\n[표현] 강조 표현: {len(self.content_patterns['강조표현'])}개")
        for expr in self.content_patterns['강조표현'][:5]:
            print(f"  - {expr[:80]}...")

        print(f"\n[데이터] 숫자 데이터: {len(self.content_patterns['숫자데이터'])}개")
        for data in self.content_patterns['숫자데이터'][:5]:
            print(f"  - {data[:80]}...")

        print(f"\n[사례] 성공 사례: {len(self.content_patterns['성공사례'])}개")
        for case in self.content_patterns['성공사례'][:3]:
            print(f"  - {case[:80]}...")

    def save_json(self, output_path):
        """학습 결과를 JSON으로 저장"""
        output_file = Path(output_path)

        learning_data = {
            "메타정보": {
                "원본파일": self.pptx_path.name,
                "총슬라이드수": len(self.prs.slides),
            },
            "학습내용": self.content_patterns
        }

        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(learning_data, f, ensure_ascii=False, indent=2)

        print(f"\n[OK] 학습 데이터 저장: {output_file}")
        return output_file

    def generate_writing_guide(self, output_path):
        """작성 가이드 Markdown 생성"""
        md_content = f"""# {self.pptx_path.stem} 작성 가이드

## 📋 섹션 구조

"""
        for section in self.content_patterns['섹션구조']:
            md_content += f"### {section['제목']}\n\n"
            if section['하위슬라이드']:
                for sub in section['하위슬라이드'][:5]:
                    md_content += f"- {sub['제목']}\n"
                md_content += "\n"

        md_content += """
## ✍️ 문장 스타일

### 강조 표현 패턴
"""
        for expr in self.content_patterns['강조표현'][:10]:
            md_content += f"- {expr}\n"

        md_content += """

### 숫자 활용
"""
        for data in self.content_patterns['숫자데이터'][:10]:
            md_content += f"- {data}\n"

        md_content += """

## 🏆 성공 사례 표현

"""
        for case in self.content_patterns['성공사례'][:10]:
            md_content += f"- {case}\n"

        md_content += """

## 🎯 AI 작성 시 주의사항

1. **섹션 순서 준수**
   - 위 섹션 구조를 그대로 유지

2. **강조 표현 활용**
   - 위 강조 표현 패턴을 참고하여 유사한 톤으로 작성

3. **숫자 근거**
   - 성과 지표는 반드시 숫자와 함께 제시

4. **성공 사례 스타일**
   - 위 성공 사례 표현을 참고하여 구체적으로 작성
"""

        output_file = Path(output_path)
        output_file.write_text(md_content, encoding='utf-8')

        print(f"작성 가이드 저장: {output_file}")
        return output_file


def main():
    if len(sys.argv) < 2:
        print("사용법: python proposal_content_learner.py <PPTX파일> [템플릿JSON] [출력폴더]")
        print("예시: python proposal_content_learner.py 리멘제안서.pptx template.json learned/")
        sys.exit(1)

    pptx_path = sys.argv[1]
    template_json = sys.argv[2] if len(sys.argv) > 2 else None
    output_dir = Path(sys.argv[3]) if len(sys.argv) > 3 else Path('.')
    output_dir.mkdir(exist_ok=True)

    learner = ProposalContentLearner(pptx_path, template_json)
    learner.learn_all()

    # JSON 저장
    stem = Path(pptx_path).stem
    json_path = output_dir / f"{stem}_content_patterns.json"
    learner.save_json(json_path)

    # 작성 가이드 생성
    guide_path = output_dir / f"{stem}_writing_guide.md"
    learner.generate_writing_guide(guide_path)

    print("\n완료! 다음 파일 생성:")
    print(f"  1. {json_path} (학습 데이터)")
    print(f"  2. {guide_path} (작성 가이드)")


if __name__ == '__main__':
    main()
