"""
PPT 완전 분석기
기존 PPT를 슬라이드별로 완전히 분석하여 재현 가능한 JSON 레시피 생성
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches


class PPTAnalyzer:
    def __init__(self, pptx_path):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(pptx_path))

        self.analysis = {
            "메타정보": {
                "원본파일": self.pptx_path.name,
                "총슬라이드수": len(self.prs.slides),
                "슬라이드크기": {
                    "width": self.prs.slide_width,
                    "height": self.prs.slide_height
                }
            },
            "슬라이드": []
        }

    def analyze_all(self):
        """전체 PPT 분석"""
        print(f"[분석] PPT 분석 시작: {self.pptx_path.name}")
        print(f"총 {len(self.prs.slides)} 슬라이드")
        print("=" * 80)

        for slide_num, slide in enumerate(self.prs.slides, 1):
            print(f"슬라이드 {slide_num} 분석 중...", end='\r')
            slide_data = self._analyze_slide(slide, slide_num)
            self.analysis["슬라이드"].append(slide_data)

        print(f"\n[OK] 분석 완료!")
        return self.analysis

    def _analyze_slide(self, slide, slide_num):
        """단일 슬라이드 분석"""
        slide_data = {
            "슬라이드번호": slide_num,
            "레이아웃": self._get_layout_name(slide),
            "요소": []
        }

        # 모든 shape 분석
        for shape_idx, shape in enumerate(slide.shapes):
            element = self._analyze_shape(shape, shape_idx)
            if element:
                slide_data["요소"].append(element)

        return slide_data

    def _get_layout_name(self, slide):
        """슬라이드 레이아웃 이름 가져오기"""
        try:
            return slide.slide_layout.name
        except:
            return "Unknown"

    def _analyze_shape(self, shape, shape_idx):
        """Shape 분석"""
        element = {
            "순서": shape_idx,
            "타입": self._get_shape_type(shape),
            "위치": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
        }

        # 텍스트가 있는 shape
        if shape.has_text_frame:
            element["텍스트"] = self._analyze_text_frame(shape.text_frame)

        # 표
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            element["표"] = self._analyze_table(shape.table)

        # 그림
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["이미지"] = {
                "설명": shape.name,
                # 이미지 데이터는 나중에 필요하면 추가
            }

        # 제목 여부
        if shape == shape.part.slide.shapes.title:
            element["제목여부"] = True

        return element

    def _get_shape_type(self, shape):
        """Shape 타입 문자열로 변환"""
        type_map = {
            MSO_SHAPE_TYPE.TEXT_BOX: "텍스트박스",
            MSO_SHAPE_TYPE.PICTURE: "이미지",
            MSO_SHAPE_TYPE.TABLE: "표",
            MSO_SHAPE_TYPE.PLACEHOLDER: "플레이스홀더",
            MSO_SHAPE_TYPE.AUTO_SHAPE: "도형",
            MSO_SHAPE_TYPE.GROUP: "그룹",
        }
        return type_map.get(shape.shape_type, f"기타({shape.shape_type})")

    def _analyze_text_frame(self, text_frame):
        """텍스트 프레임 분석"""
        paragraphs_data = []

        for para in text_frame.paragraphs:
            para_data = {
                "텍스트": para.text,
                "레벨": para.level,
                "정렬": str(para.alignment) if para.alignment else None,
                "runs": []
            }

            # Run별 스타일 분석
            for run in para.runs:
                run_data = {
                    "텍스트": run.text,
                    "폰트": {
                        "이름": run.font.name,
                        "크기": run.font.size.pt if run.font.size else None,
                        "굵게": run.font.bold,
                        "기울임": run.font.italic,
                        "밑줄": run.font.underline,
                    }
                }

                # 색상 (있으면)
                try:
                    if run.font.color and run.font.color.rgb:
                        run_data["폰트"]["색상"] = str(run.font.color.rgb)
                except (AttributeError, TypeError):
                    pass  # SchemeColor 등 rgb 속성이 없는 경우

                para_data["runs"].append(run_data)

            paragraphs_data.append(para_data)

        return {
            "단락": paragraphs_data
        }

    def _analyze_table(self, table):
        """표 분석"""
        table_data = {
            "행수": len(table.rows),
            "열수": len(table.columns),
            "셀": []
        }

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_data = {
                    "행": row_idx,
                    "열": col_idx,
                    "텍스트": cell.text,
                    "병합여부": cell.is_merge_origin or cell.is_spanned
                }
                table_data["셀"].append(cell_data)

        return table_data

    def save_json(self, output_path):
        """분석 결과를 JSON으로 저장"""
        output_file = Path(output_path)

        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(self.analysis, f, ensure_ascii=False, indent=2)

        print(f"\n[OK] 분석 데이터 저장: {output_file}")
        return output_file

    def print_summary(self):
        """분석 결과 요약 출력"""
        print("\n" + "=" * 80)
        print("[요약] 분석 결과 요약")
        print("=" * 80)

        print(f"\n총 슬라이드: {self.analysis['메타정보']['총슬라이드수']}개")

        # 슬라이드별 요소 수
        for slide in self.analysis["슬라이드"][:5]:  # 처음 5개만
            num_elements = len(slide["요소"])
            print(f"  슬라이드 {slide['슬라이드번호']}: {num_elements}개 요소")

        # 타입별 통계
        type_counts = {}
        for slide in self.analysis["슬라이드"]:
            for element in slide["요소"]:
                element_type = element["타입"]
                type_counts[element_type] = type_counts.get(element_type, 0) + 1

        print(f"\n요소 타입별 통계:")
        for element_type, count in sorted(type_counts.items(), key=lambda x: -x[1]):
            print(f"  - {element_type}: {count}개")


def main():
    if len(sys.argv) < 2:
        print("사용법: python ppt_analyzer.py <PPTX파일> [출력JSON]")
        print("예시: python ppt_analyzer.py 제안서.pptx analysis.json")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_json = sys.argv[2] if len(sys.argv) > 2 else None

    if not output_json:
        # 기본 출력 파일명: 원본명_analysis.json
        stem = Path(pptx_path).stem
        output_json = f"{stem}_analysis.json"

    analyzer = PPTAnalyzer(pptx_path)
    analyzer.analyze_all()
    analyzer.print_summary()
    analyzer.save_json(output_json)

    print("\n완료! 분석 데이터가 저장되었습니다.")
    print(f"파일: {output_json}")


if __name__ == '__main__':
    main()
