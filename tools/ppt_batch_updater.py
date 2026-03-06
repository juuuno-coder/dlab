"""
PPT 배치 업데이터
JSON 설정 파일로 여러 치환 규칙을 한 번에 처리
"""

import sys
import json
from pathlib import Path
from pptx import Presentation


class PPTBatchUpdater:
    def __init__(self, pptx_path, batch_config_path):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(pptx_path))

        # 배치 설정 로드
        with open(batch_config_path, 'r', encoding='utf-8') as f:
            self.config = json.load(f)

        print(f"[배치 업데이트] PPT 로드: {self.pptx_path.name}")
        print(f"총 {len(self.prs.slides)} 슬라이드")
        print(f"배치 설정 로드: {Path(batch_config_path).name}")

        # 치환 규칙 준비
        self.replacements = self.config.get('replacements', {})
        print(f"치환 규칙: {len(self.replacements)}개")

    def preview(self):
        """
        배치 업데이트 미리보기
        """
        print("\n[미리보기] 다음 치환이 수행됩니다:")
        print("=" * 80)

        for old_value, new_value in self.replacements.items():
            # 각 값이 몇 번 발견되는지 확인
            count = self._count_occurrences(old_value)
            print(f"  '{old_value}' → '{new_value}' ({count}건)")

        total = sum(self._count_occurrences(old) for old in self.replacements.keys())
        print(f"\n총 {total}개 항목이 변경됩니다.")

    def _count_occurrences(self, search_text):
        """특정 텍스트가 몇 번 나오는지 카운트"""
        count = 0

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if search_text in run.text:
                            count += 1

        return count

    def execute(self):
        """
        배치 업데이트 실행
        충돌 방지 알고리즘 적용:
        - 긴 문자열부터 치환 (A→B, AB→C일 때 AB를 먼저 처리)
        - 임시 플레이스홀더 사용 (A→B, B→C 충돌 방지)
        """
        print("\n[배치 실행]")
        print("=" * 80)

        # 1단계: 긴 문자열부터 정렬
        sorted_replacements = sorted(
            self.replacements.items(),
            key=lambda x: len(x[0]),
            reverse=True
        )

        # 2단계: 임시 플레이스홀더 생성
        temp_map = {}
        for idx, (old_value, new_value) in enumerate(sorted_replacements):
            temp_placeholder = f"___TEMP_{idx}___"
            temp_map[old_value] = (temp_placeholder, new_value)

        # 3단계: 먼저 임시 플레이스홀더로 치환
        print("\n[1단계] 임시 플레이스홀더로 치환 중...")
        for old_value, (placeholder, _) in temp_map.items():
            count = self._replace_text(old_value, placeholder)
            if count > 0:
                print(f"  '{old_value}' → '{placeholder}' ({count}건)")

        # 4단계: 플레이스홀더를 실제 값으로 치환
        print("\n[2단계] 최종 값으로 치환 중...")
        total_replaced = 0
        for old_value, (placeholder, new_value) in temp_map.items():
            count = self._replace_text(placeholder, new_value)
            if count > 0:
                print(f"  '{placeholder}' → '{new_value}' ({count}건)")
                total_replaced += count

        print(f"\n[OK] 총 {total_replaced}개 항목 변경 완료")
        return total_replaced

    def _replace_text(self, old_text, new_text):
        """텍스트 치환"""
        replaced_count = 0

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            replaced_count += 1

        return replaced_count

    def save(self, output_path=None):
        """저장"""
        if output_path is None:
            # 원본 파일 이름에 _batch_updated 추가
            stem = self.pptx_path.stem
            suffix = self.pptx_path.suffix
            output_path = self.pptx_path.parent / f"{stem}_batch_updated{suffix}"

        self.prs.save(str(output_path))
        print(f"\n[저장 완료] {output_path}")
        return output_path


def main():
    if len(sys.argv) < 3:
        print("사용법: python ppt_batch_updater.py <PPTX파일> <배치설정JSON>")
        print("\n배치 설정 JSON 형식:")
        print('{')
        print('  "replacements": {')
        print('    "2026년 제6회 부산 봄꽃 전시회": "2026년 광주 가을 축제",')
        print('    "110,000,000원": "80,000,000원",')
        print('    "부산시민공원": "광주문화예술회관"')
        print('  }')
        print('}')
        print("\n예시:")
        print("  python ppt_batch_updater.py proposal.pptx batch_config.json")
        sys.exit(1)

    pptx_path = sys.argv[1]
    batch_config_path = sys.argv[2]

    try:
        updater = PPTBatchUpdater(pptx_path, batch_config_path)

        # 미리보기
        updater.preview()

        # 확인
        if '--yes' not in sys.argv and '-y' not in sys.argv:
            response = input("\n계속하시겠습니까? (y/n): ")
            if response.lower() != 'y':
                print("취소되었습니다.")
                sys.exit(0)

        # 실행
        updated_count = updater.execute()

        if updated_count > 0:
            # 저장
            output_path = updater.save()

            print("\n" + "=" * 80)
            print("[완료] PPT 배치 업데이트 완료!")
            print("=" * 80)
            print(f"\n업데이트된 파일: {output_path}")

    except Exception as e:
        print(f"\n[오류] 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
