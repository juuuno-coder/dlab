"""
PPT 템플릿 분석기 - AI가 학습할 수 있도록 구조 추출
python-pptx로 슬라이드별 패턴 분석
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from collections import defaultdict


class PPTTemplateAnalyzer:
    def __init__(self, pptx_path):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(pptx_path))
        self.template_data = {
            'metadata': {
                'file_name': self.pptx_path.name,
                'total_slides': len(self.prs.slides),
                'slide_width': self.prs.slide_width,
                'slide_height': self.prs.slide_height
            },
            'slide_patterns': [],
            'text_styles': defaultdict(list),
            'design_system': {
                'colors': set(),
                'fonts': set(),
                'layout_types': defaultdict(int)
            }
        }

    def classify_slide_type(self, slide, slide_num):
        """슬라이드 타입 분류"""
        shapes_count = len(slide.shapes)
        has_title = any(shape.has_text_frame and
                       getattr(shape, 'is_placeholder', False) and
                       shape.placeholder_format.type == 1
                       for shape in slide.shapes)

        # 텍스트 양 계산
        total_text_length = sum(len(shape.text) for shape in slide.shapes if hasattr(shape, 'text'))

        # 이미지/표 개수
        has_image = any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = PICTURE
        has_table = any(shape.shape_type == 19 for shape in slide.shapes)  # 19 = TABLE

        # 슬라이드 타입 판단
        if slide_num == 1:
            return 'cover'
        elif total_text_length < 50 and shapes_count < 5:
            return 'divider'
        elif has_table:
            return 'table'
        elif has_image:
            return 'image-content'
        elif total_text_length > 300:
            return 'text-heavy'
        else:
            return 'content'

    def extract_text_style(self, run):
        """텍스트 런(Run)의 스타일 정보 추출"""
        try:
            return {
                'font_name': run.font.name if run.font.name else 'default',
                'font_size': run.font.size.pt if run.font.size else None,
                'bold': run.font.bold if run.font.bold is not None else False,
                'italic': run.font.italic if run.font.italic is not None else False,
                'color': self._get_color_rgb(run.font.color) if run.font.color else None
            }
        except:
            return {}

    def _get_color_rgb(self, color):
        """색상 RGB 추출"""
        try:
            if color.type == 1:  # RGB 색상
                return f"#{color.rgb:06x}"
        except:
            pass
        return None

    def analyze_slide(self, slide, slide_num):
        """단일 슬라이드 분석"""
        slide_type = self.classify_slide_type(slide, slide_num)

        slide_data = {
            'slide_number': slide_num,
            'slide_type': slide_type,
            'shapes': []
        }

        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {
                'shape_id': shape_idx,
                'shape_type': str(shape.shape_type),
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }

            # 텍스트가 있는 경우
            if hasattr(shape, 'text_frame'):
                shape_info['has_text'] = True
                shape_info['text_content'] = shape.text[:200]  # 처음 200자
                shape_info['paragraphs'] = []

                for para in shape.text_frame.paragraphs:
                    para_info = {
                        'text': para.text[:100],
                        'level': para.level,
                        'alignment': str(para.alignment) if para.alignment else 'LEFT',
                        'runs': []
                    }

                    for run in para.runs:
                        run_style = self.extract_text_style(run)
                        para_info['runs'].append(run_style)

                        # 폰트/색상 수집
                        if run_style.get('font_name'):
                            self.template_data['design_system']['fonts'].add(run_style['font_name'])
                        if run_style.get('color'):
                            self.template_data['design_system']['colors'].add(run_style['color'])

                    shape_info['paragraphs'].append(para_info)

            # 표가 있는 경우
            if shape.shape_type == 19:  # TABLE
                shape_info['table'] = {
                    'rows': len(shape.table.rows),
                    'columns': len(shape.table.columns)
                }

            slide_data['shapes'].append(shape_info)

        # 레이아웃 타입 카운트
        self.template_data['design_system']['layout_types'][slide_type] += 1

        return slide_data

    def generate_template_summary(self):
        """템플릿 요약 생성"""
        summary = {
            '총 슬라이드 수': self.template_data['metadata']['total_slides'],
            '슬라이드 타입별 분포': dict(self.template_data['design_system']['layout_types']),
            '사용된 폰트': list(self.template_data['design_system']['fonts']),
            '사용된 색상': list(self.template_data['design_system']['colors']),
        }

        return summary

    def analyze_all(self):
        """전체 PPT 분석"""
        print(f"PPT 템플릿 분석 시작: {self.pptx_path.name}")
        print(f"총 {len(self.prs.slides)} 슬라이드")
        print("=" * 80)

        for slide_num, slide in enumerate(self.prs.slides, 1):
            print(f"슬라이드 {slide_num} 분석 중...", end='\r')
            slide_data = self.analyze_slide(slide, slide_num)
            self.template_data['slide_patterns'].append(slide_data)

        print(f"\n분석 완료!")

        # 요약 출력
        summary = self.generate_template_summary()
        print("\n" + "=" * 80)
        print("템플릿 요약:")
        print("=" * 80)
        for key, value in summary.items():
            print(f"{key}: {value}")

        return self.template_data

    def save_template(self, output_path):
        """템플릿을 JSON으로 저장"""
        # set을 list로 변환
        self.template_data['design_system']['fonts'] = list(self.template_data['design_system']['fonts'])
        self.template_data['design_system']['colors'] = list(self.template_data['design_system']['colors'])
        self.template_data['design_system']['layout_types'] = dict(self.template_data['design_system']['layout_types'])

        output_file = Path(output_path)
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(self.template_data, f, ensure_ascii=False, indent=2)

        print(f"\n템플릿 저장: {output_file}")
        return output_file

    def generate_markdown_guide(self, output_path):
        """템플릿 가이드를 Markdown으로 생성"""
        summary = self.generate_template_summary()

        md_content = f"""# {self.pptx_path.stem} 템플릿 가이드

## 📊 기본 정보
- **파일명**: {self.pptx_path.name}
- **총 슬라이드**: {summary['총 슬라이드 수']}개
- **슬라이드 크기**: {self.template_data['metadata']['slide_width'] / 914400:.0f} x {self.template_data['metadata']['slide_height'] / 914400:.0f} inches

## 🎨 디자인 시스템

### 사용된 폰트
"""
        for font in summary['사용된 폰트']:
            md_content += f"- {font}\n"

        md_content += "\n### 사용된 색상\n"
        for color in summary['사용된 색상'][:10]:  # 상위 10개
            md_content += f"- `{color}`\n"

        md_content += "\n## 📑 슬라이드 타입별 구조\n\n"
        for slide_type, count in summary['슬라이드 타입별 분포'].items():
            md_content += f"### {slide_type} ({count}개)\n\n"

            # 해당 타입의 첫 번째 슬라이드 예시
            example = next((s for s in self.template_data['slide_patterns']
                          if s['slide_type'] == slide_type), None)

            if example:
                md_content += f"**예시: 슬라이드 {example['slide_number']}**\n\n"
                md_content += f"- 도형 개수: {len(example['shapes'])}개\n"

                # 텍스트가 있는 도형만
                text_shapes = [s for s in example['shapes'] if s.get('has_text')]
                if text_shapes:
                    md_content += f"- 텍스트 블록: {len(text_shapes)}개\n"

                    # 첫 번째 텍스트 샘플
                    first_text = text_shapes[0].get('text_content', '')
                    if first_text:
                        md_content += f"- 텍스트 샘플: \"{first_text[:50]}...\"\n"

            md_content += "\n"

        md_content += """
## 🔄 AI 재생성 가이드

이 템플릿을 기반으로 새로운 제안서를 생성할 때:

1. **슬라이드 순서 유지**
   - 표지 → 목차 → 본문 → 마무리 순서 준수

2. **디자인 일관성**
   - 위 폰트 목록의 폰트만 사용
   - 위 색상 팔레트 내에서만 사용

3. **텍스트 스타일**
   - 제목: [분석 결과에서 추출된 폰트/크기]
   - 본문: [분석 결과에서 추출된 폰트/크기]
   - 강조: 볼드 + 색상 변경

4. **레이아웃 패턴**
   - 각 슬라이드 타입별로 위 구조를 따름
"""

        output_file = Path(output_path)
        output_file.write_text(md_content, encoding='utf-8')

        print(f"템플릿 가이드 저장: {output_file}")
        return output_file


def main():
    if len(sys.argv) < 2:
        print("사용법: python ppt_template_analyzer.py <PPTX파일> [출력폴더]")
        print("예시: python ppt_template_analyzer.py 리멘제안서.pptx templates/")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else Path('.')
    output_dir.mkdir(exist_ok=True)

    analyzer = PPTTemplateAnalyzer(pptx_path)
    analyzer.analyze_all()

    # JSON 템플릿 저장
    stem = Path(pptx_path).stem
    json_path = output_dir / f"{stem}_template.json"
    analyzer.save_template(json_path)

    # Markdown 가이드 생성
    md_path = output_dir / f"{stem}_guide.md"
    analyzer.generate_markdown_guide(md_path)

    print("\n완료! 다음 파일 생성:")
    print(f"  1. {json_path} (구조화된 템플릿 데이터)")
    print(f"  2. {md_path} (AI용 가이드 문서)")


if __name__ == '__main__':
    main()
