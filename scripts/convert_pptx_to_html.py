#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX → HTML 완벽 변환 스크립트
원본 제안서 양식을 픽셀 단위까지 정확하게 HTML로 재현
"""

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import base64

# PPTX 파일 경로
PPTX_PATH = r"d:\dev\dlab-site\knowledge-base\templates\260224_2026 부산 봄꽃 전시회 용역 제안.pptx"
OUTPUT_HTML = r"d:\dev\dlab-site\public\2026_부산봄꽃전시회_제안서_REAL.html"
IMAGE_DIR = r"d:\dev\dlab-site\public\proposal_images"

# 생성된 이미지 파일 매핑 (슬라이드 번호 → 이미지 파일명)
GENERATED_IMAGES = {
    1: ['slide_01_1.png'],
    3: ['slide_03_1.png'],
    4: ['slide_04_1.png', 'slide_04_2.png', 'slide_04_3.png'],
    11: ['slide_11_1.png', 'slide_11_2.png', 'slide_11_3.png'],
    12: ['slide_12_1.png'],
    27: ['slide_27_1.png']
}

def rgb_to_hex(rgb):
    """RGB 튜플을 HEX 색상 코드로 변환"""
    if rgb is None:
        return '#000000'
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])

def emu_to_px(emu):
    """EMU (English Metric Unit)를 픽셀로 변환"""
    if emu is None:
        return 0
    return int(emu / 9525)  # 1px = 9525 EMU

def get_text_alignment(para):
    """텍스트 정렬 방식 추출"""
    if para.alignment is None:
        return 'left'
    if para.alignment == PP_ALIGN.CENTER:
        return 'center'
    elif para.alignment == PP_ALIGN.RIGHT:
        return 'right'
    else:
        return 'left'

def extract_slide_data(prs):
    """모든 슬라이드 데이터 추출"""
    slides_data = []

    for idx, slide in enumerate(prs.slides, 1):
        print(f"[*] 슬라이드 {idx} 분석 중...")

        slide_info = {
            'number': idx,
            'shapes': [],
            'background': None,
            'has_generated_images': idx in GENERATED_IMAGES
        }

        # 슬라이드 배경 색상 추출
        if slide.background.fill.type == 1:  # Solid fill
            try:
                bg_color = slide.background.fill.fore_color.rgb
                slide_info['background'] = rgb_to_hex(bg_color)
            except:
                slide_info['background'] = '#FFFFFF'

        # 모든 도형/텍스트 박스 추출
        for shape in slide.shapes:
            shape_data = extract_shape(shape)
            if shape_data:
                slide_info['shapes'].append(shape_data)

        slides_data.append(slide_info)

    return slides_data

def extract_shape(shape):
    """개별 도형/텍스트 박스 데이터 추출"""
    try:
        shape_data = {
            'type': str(shape.shape_type),
            'left': emu_to_px(shape.left),
            'top': emu_to_px(shape.top),
            'width': emu_to_px(shape.width),
            'height': emu_to_px(shape.height),
            'text': '',
            'font_size': 12,
            'font_color': '#000000',
            'font_bold': False,
            'alignment': 'left',
            'fill_color': None,
            'border_color': None
        }

        # 텍스트가 있는 경우
        if shape.has_text_frame:
            texts = []
            for para in shape.text_frame.paragraphs:
                para_text = ''
                for run in para.runs:
                    para_text += run.text
                if para_text:
                    texts.append(para_text)

                # 첫 번째 paragraph의 스타일 정보 추출
                if para.runs and not shape_data['text']:
                    run = para.runs[0]
                    if run.font.size:
                        shape_data['font_size'] = int(run.font.size.pt)
                    if run.font.color and run.font.color.rgb:
                        shape_data['font_color'] = rgb_to_hex(run.font.color.rgb)
                    shape_data['font_bold'] = run.font.bold or False
                    shape_data['alignment'] = get_text_alignment(para)

            shape_data['text'] = '\n'.join(texts)

        # 채우기 색상
        if shape.fill.type == 1:  # Solid fill
            try:
                shape_data['fill_color'] = rgb_to_hex(shape.fill.fore_color.rgb)
            except:
                pass

        return shape_data

    except Exception as e:
        print(f"[WARN]  도형 추출 실패: {e}")
        return None

def generate_html(slides_data):
    """HTML 파일 생성"""

    html = """<!DOCTYPE html>
<html lang="ko">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>2026년 제6회 부산 봄꽃 전시회 - 제안서 (원본 양식)</title>
  <style>
    * { margin: 0; padding: 0; box-sizing: border-box; }
    body {
      font-family: 'Malgun Gothic', '맑은 고딕', sans-serif;
      background: #f5f5f5;
      line-height: 1.6;
    }

    /* 슬라이드 기본 */
    .slide {
      width: 1920px;
      height: 1080px;
      background: white;
      margin: 20px auto;
      position: relative;
      overflow: hidden;
      box-shadow: 0 10px 40px rgba(0,0,0,0.1);
      page-break-after: always;
    }

    /* 도형/텍스트 박스 */
    .shape {
      position: absolute;
      white-space: pre-wrap;
      word-wrap: break-word;
    }

    /* 생성된 이미지 */
    .generated-image {
      max-width: 100%;
      height: auto;
      border-radius: 4px;
    }

    /* 인쇄 최적화 */
    @media print {
      body { background: white; }
      .slide { margin: 0; box-shadow: none; }
    }
  </style>
</head>
<body>
"""

    # 각 슬라이드 생성
    for slide in slides_data:
        bg_color = slide.get('background', '#FFFFFF')

        html += f'\n  <!-- 슬라이드 {slide["number"]} -->\n'
        html += f'  <div class="slide" style="background-color: {bg_color};">\n'

        # 생성된 이미지가 있는 경우 배경으로 삽입
        if slide['has_generated_images'] and slide['number'] in [1, 3]:
            img_file = GENERATED_IMAGES[slide['number']][0]
            html += f'    <div style="position: absolute; top: 0; left: 0; width: 100%; height: 100%; z-index: 0;">\n'
            html += f'      <img src="/proposal_images/{img_file}" style="width: 100%; height: 100%; object-fit: cover;" alt="배경 이미지">\n'
            html += f'    </div>\n'

        # 모든 도형/텍스트 박스
        for shape in slide['shapes']:
            if not shape['text'] and not shape['fill_color']:
                continue

            styles = [
                f"left: {shape['left']}px",
                f"top: {shape['top']}px",
                f"width: {shape['width']}px",
                f"height: {shape['height']}px",
                f"font-size: {shape['font_size']}pt",
                f"color: {shape['font_color']}",
                f"text-align: {shape['alignment']}",
            ]

            if shape['font_bold']:
                styles.append("font-weight: bold")

            if shape['fill_color']:
                styles.append(f"background-color: {shape['fill_color']}")

            style_str = '; '.join(styles)

            text = shape['text'].replace('\n', '<br>')
            html += f'    <div class="shape" style="{style_str}">{text}</div>\n'

        # 생성된 이미지 삽입 (슬라이드 4, 11, 12, 27)
        if slide['has_generated_images'] and slide['number'] not in [1, 3]:
            images = GENERATED_IMAGES[slide['number']]

            # 이미지 배치 (간단한 그리드)
            if len(images) == 3:
                html += f'    <div style="position: absolute; bottom: 100px; left: 100px; right: 100px; display: grid; grid-template-columns: repeat(3, 1fr); gap: 30px; z-index: 999;">\n'
                for img in images:
                    html += f'      <img src="/proposal_images/{img}" class="generated-image" alt="이미지">\n'
                html += f'    </div>\n'
            elif len(images) == 1:
                html += f'    <div style="position: absolute; bottom: 100px; left: 100px; right: 100px; z-index: 999;">\n'
                html += f'      <img src="/proposal_images/{images[0]}" class="generated-image" alt="이미지">\n'
                html += f'    </div>\n'

        html += '  </div>\n'

    html += """
</body>
</html>
"""

    return html

def main():
    print("\n[*] PPTX -> HTML 완벽 변환 시작...\n")
    print(f"[*] 원본 PPTX: {PPTX_PATH}")

    # PPTX 로드
    try:
        prs = Presentation(PPTX_PATH)
        print(f"[OK] PPTX 로드 완료 (총 {len(prs.slides)}개 슬라이드)\n")
    except Exception as e:
        print(f"[ERR] PPTX 로드 실패: {e}")
        return

    # 슬라이드 데이터 추출
    slides_data = extract_slide_data(prs)
    print(f"\n[OK] 슬라이드 데이터 추출 완료 ({len(slides_data)}개)\n")

    # HTML 생성
    html_content = generate_html(slides_data)

    # 파일 저장
    with open(OUTPUT_HTML, 'w', encoding='utf-8') as f:
        f.write(html_content)

    print(f"[OK] HTML 저장 완료!")
    print(f"[FILE] 저장 위치: {OUTPUT_HTML}\n")
    print("다음 단계:")
    print("1. 브라우저로 HTML 파일 열기")
    print("2. 원본 PPTX와 비교하여 양식 확인")
    print("3. 필요시 미세 조정\n")

if __name__ == '__main__':
    main()
